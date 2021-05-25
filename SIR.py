from tools import *
##### Don't need to run - one time task of reading the dump and creating a corpus
rerun=False
if rerun:
    loc ="/Users/akhilesh/Desktop/Presentations/"
    universal_doc=doc_creator(loc) 
    corpus= cleanit(list(universal_doc.values()))
    vectorizer=vectorizer_creation(corpus)
    slide_names=list(universal_doc.keys())
    corp_vector=list(vector(corpus,vectorizer))
    f = open('/Users/akhilesh/Desktop/Presentations/qER for pharma (tPA).pptx', "rb")
    standard=Presentation(f)
    pickle.dump(universal_doc, open( "/Users/akhilesh/Desktop/universal_doc.pickle", "wb" ))
    pickle.dump(corpus, open( "/Users/akhilesh/Desktop/corpus.pickle", "wb" ))
    
    
path11='/Users/akhilesh/Desktop/'#input('Please enter the directory of the folder with pickle files:')
universal_doc = pickle.load( open( path11+"universal_doc.pickle", "rb" ))
corpus=pickle.load( open( path11+"corpus.pickle", "rb" ))
vectorizer=vectorizer_creation(corpus)
slide_names=list(universal_doc.keys())
corp_vector=list(vector(corpus,vectorizer))
tf_vectorizer=tf_vectorizer_creation(corpus)
corp_tf_vector=tf_vector(corpus,tf_vectorizer)

def get_slide(query):
    def tf_notsameslide(slide,final_vector):
        if len(final_vector)==0:
            return True
        sim=[]
        for i in range(len(final_vector)):
            sim.append(1-cosine(final_vector[i],slide))
        #print(sim)
        ind=sim.index(max(sim))
        if max(sim)>0.95:
            #print("in",'\n')
            if np.count_nonzero(final_vector[ind])>=np.count_nonzero(slide):
                return False
            else:
                del final_vectors[ind]
                del final_result[ind]
                return True
        else:
            return True
    key=query
    final_result=[]
    final_vectors=[]
    key_vector=list(vector(cleanit([key])*len(corpus),vectorizer))
    result=np.einsum('ij, ij->i', key_vector, corp_vector)#np.multiply(key_vector, corp_vector).sum(1)
    similarity=result/sum(key_vector[0])
    for i in range(len(result)):
        if similarity[i]>=0.75:
            if tf_notsameslide(corp_tf_vector[i],final_vectors):
                final_result.append(slide_names[i])
                final_vectors.append(corp_tf_vector[i])
    print('\n'.join(final_result))
    #creates ppts and saves them to a location provided by the user
    from pptx import Presentation
    source_dir='/Users/akhilesh/Desktop/Presentations/'
    target_dir=source_dir#input("Please enter the directory where the ppt needs to be saved : ")
    searchresult1=Presentation()
    searchresult2=Presentation()
    print('Generating PPT',end='')
    for i in range(len(final_result)):
        source=source_dir+' '.join(final_result[i].split()[:-2])
        f = open(source, "rb")
        prs=Presentation(f)
        slide=final_result[i].split()[-1]
        copy_slide_from_external_prs(prs,int(slide), searchresult1,searchresult2)
        print('.',end='')
    searchresult1.save(target_dir+'search_result1.pptx')
    if len(searchresult2.slides)>0:
        searchresult2.save(target_dir+'search_result2.pptx')
    return('Your PPTs are saved at '+target_dir)

from flask import Flask
from flask_restplus import Api, Resource, reqparse
from werkzeug.utils import cached_property
app = Flask(__name__)
api = Api(app, version='1.0', title='SIR - Slide Identifier and Retriever', description='Identifies slides with keywords and retrieves them')
ns = api.namespace('SIR', description='')
parser = reqparse.RequestParser()
parser.add_argument('keywords', help='Specify the keywords:')
@ns.route('/Find slides/')
class slide_finder(Resource):
    @api.doc(parser=parser)
    def get(self):        
        args = parser.parse_args()
        keywords = args['keywords']
        output=get_slide(keywords)
        return output
    
if __name__ == '__main__':
    app.run()

print('Server is running')