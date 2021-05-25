import copy
import os
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
import glob
import os
from scipy.spatial.distance import cosine
import pickle
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.stem import PorterStemmer
from cleantext import clean
from gensim.parsing.preprocessing import remove_stopwords
from keras.preprocessing.text import Tokenizer
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np
def doc_creator(path):
    universal_doc={}
    for eachfile in glob.glob(os.path.join(path, '*.pptx')):
        filename=eachfile.replace(path,'')
        if filename[0:2]=='~$':
            continue
        f = open(eachfile, "rb")
        prs=Presentation(f)
        sd=0
        for slide in prs.slides:
            sd+=1
            universal_doc[filename+" Slide: "+str(sd)]=''
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    universal_doc[filename+" Slide: "+str(sd)]+=shape.text+" "
    return (universal_doc)  

def remove_dup(s):
    l = s.split()
    k = []
    for i in l:

        # If condition is used to store unique string 
        # in another list 'k' 
        if (s.count(i)>1 and (i not in k)or s.count(i)==1):
            k.append(i)
    return ' '.join(k)

def stemSentence(sentence):
    porter=PorterStemmer()
    token_words=list(sentence.split())
    token_words
    stem_sentence=[]
    for word in token_words:
        stem_sentence.append(porter.stem(word))
        stem_sentence.append(" ")
    return "".join(stem_sentence)

def cleanit(file):
    for i in range(len(file)):
        file[i]=remove_stopwords(file[i])
        file[i]=clean(file[i],no_line_breaks=True,no_numbers=True,no_digits=True,no_currency_symbols=True,no_punct=True, replace_with_punct=" ",replace_with_number="",
            replace_with_digit="")
        file[i]=stemSentence(file[i])
        file[i]=remove_dup(file[i])
    return file  
   
def vectorizer_creation(corpus):
    tokenizer = Tokenizer()
    tokenizer.fit_on_texts(corpus)
    return tokenizer
     
def vector(text,vectorizer):
    return vectorizer.texts_to_matrix(text,mode="count")

def tf_vectorizer_creation(corpus):
    tf_vectorizer = TfidfVectorizer()
    tf_vectorizer.fit(corpus) 
    return tf_vectorizer

def tf_vector(text,tf_vectorizer):
    return tf_vectorizer.transform(text).toarray()

def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]

def copy_slide_from_external_prs(src, idx, newPrs1,newPrs2):
    src_height=src.slide_height.cm
    src_width=src.slide_width.cm
    src_slide = src.slides[idx-1]
    f = open('/Users/akhilesh/Desktop/Presentations/qER for pharma (tPA).pptx', "rb")
    standard=Presentation(f)
    blank_slide_layout = _get_blank_slide_layout(standard)
    if src_width>40:
        newPrs=newPrs2
        newPrs.slide_width = Inches(20)
        newPrs.slide_height = Inches(11.251969)
    else:
        newPrs=newPrs1
        newPrs.slide_width = Inches(13.334646)
        newPrs.slide_height = Inches(7.5)
        
    curr_slide = newPrs.slides.add_slide(blank_slide_layout)
    imgDict = {}
    for shp in src_slide.shapes:
        if 'Picture' in shp.name:
            with open(shp.name+'.jpg', 'wb') as f:
                f.write(shp.image.blob)
                imgDict[shp.name+'.jpg'] = [shp.left, shp.top, shp.width, shp.height]
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            for s in shp.shapes:
                if 'Picture' in s.name:
                    with open(s.name+'.jpg', 'wb') as f:
                        f.write(s.image.blob)
                        curr_slide.shapes.add_picture(s.name+'.jpg',s.left, s.top, s.width, s.height)
                        os.remove(s.name+'.jpg')
                        #imgDict[s.name+'.jpg'] = [s.left, s.top, s.width, s.height]
    for k, v in imgDict.items():
        curr_slide.shapes.add_picture(k,v[0], int(v[1]), int(v[2]), int(v[3]))
        os.remove(k)
    for shp in src_slide.shapes:
        if 'Picture' not in shp.name:
            el = shp.element
            newel = copy.deepcopy(el)
            curr_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            

    
